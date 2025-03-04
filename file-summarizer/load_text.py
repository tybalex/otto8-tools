import io
import struct
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from helper import load_from_gptscript_workspace, save_to_gptscript_workspace


def extract_text_from_pdf(pdf_bytes: bytes) -> str:
    """Extracts text from a PDF file given as bytes."""
    pdf_stream = io.BytesIO(pdf_bytes)
    doc = fitz.open(stream=pdf_stream, filetype="pdf")
    
    text = []
    for page in doc:
        text.append(page.get_text("text"))
    
    return "\n".join(text)

def extract_text_from_docx(docx_bytes: bytes) -> str:
    """Extracts text from a Word (.docx) file given as bytes."""
    doc_stream = io.BytesIO(docx_bytes)
    doc = docx.Document(doc_stream)
    
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    """Extracts text from a PowerPoint (.pptx) file given as bytes."""
    ppt_stream = io.BytesIO(pptx_bytes)
    prs = Presentation(ppt_stream)
    
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    return "\n".join(text)



SUPPORTED_TEXT_FILE_TYPES = (".md", ".txt", ".markdown", ".text", ".mdx", ".mdtxt", ".mdtxtx")
SUPPORTED_DOC_FILE_TYPES = (".docx", ".pdf", ".pptx")
ALL_SUPPORTED_FILE_TYPES = SUPPORTED_TEXT_FILE_TYPES + SUPPORTED_DOC_FILE_TYPES

async def load_text_from_file(file_path: str) -> str:
    
    if not file_path.endswith(ALL_SUPPORTED_FILE_TYPES):
        raise ValueError(
            f"Error: the input file must end with one of the following file types: {ALL_SUPPORTED_FILE_TYPES}, other file types are not supported yet."
        )

    try:
        file_content = await load_from_gptscript_workspace(file_path)
    except Exception as e:
        raise ValueError(
            f"Failed to load file from GPTScript workspace file {file_path}, Error: {e}"
        )
    
    if file_path.endswith(SUPPORTED_TEXT_FILE_TYPES):
        file_content = file_content.decode("utf-8")
        return file_content
    elif file_path.endswith(SUPPORTED_DOC_FILE_TYPES):
        if file_path.endswith(".pdf"):
            return extract_text_from_pdf(file_content)
        elif file_path.endswith(".docx"):
            return extract_text_from_docx(file_content)
        elif file_path.endswith(".pptx"):
            return extract_text_from_pptx(file_content)
    else:
        raise ValueError(f"Unsupported file type: {file_path}")
