import gptscript
import os
import logging
import sys
import io
import struct
import fitz  # PyMuPDF
import docx
from pptx import Presentation


def setup_logger(name):
    """Setup a logger that writes to sys.stderr. This will eventually show up in GPTScript's debugging logs.

    Args:
        name (str): The name of the logger.

    Returns:
        logging.Logger: The logger.
    """
    # Create a logger
    logger = logging.getLogger(name)
    logger.setLevel(logging.DEBUG)  # Set the logging level

    # Create a stream handler that writes to sys.stderr
    stderr_handler = logging.StreamHandler(sys.stderr)

    # Create a log formatter
    formatter = logging.Formatter(
        "[File Summarizer Debugging Log]: %(asctime)s - %(name)s - %(levelname)s - %(message)s"
    )
    stderr_handler.setFormatter(formatter)

    # Add the handler to the logger
    logger.addHandler(stderr_handler)

    return logger


logger = setup_logger(__name__)

def _prepend_base_path(base_path: str, file_path: str):
    """
    Prepend a base path to a file path if it's not already rooted in the base path.

    Args:
        base_path (str): The base path to prepend.
        file_path (str): The file path to check and modify.

    Returns:
        str: The modified file path with the base path prepended if necessary.

    Examples:
      >>> prepend_base_path("files", "my-file.txt")
      'files/my-file.txt'

      >>> prepend_base_path("files", "files/my-file.txt")
      'files/my-file.txt'

      >>> prepend_base_path("files", "foo/my-file.txt")
      'files/foo/my-file.txt'

      >>> prepend_base_path("files", "bar/files/my-file.txt")
      'files/bar/files/my-file.txt'

      >>> prepend_base_path("files", "files/bar/files/my-file.txt")
      'files/bar/files/my-file.txt'
    """
    # Split the file path into parts for checking
    file_parts = os.path.normpath(file_path).split(os.sep)

    # Check if the base path is already at the root
    if file_parts[0] == base_path:
        return file_path

    # Prepend the base path
    return os.path.join(base_path, file_path)


# for gptscript workspace S/L, see https://github.com/gptscript-ai/py-gptscript/blob/main/gptscript/gptscript.py
async def save_to_gptscript_workspace(filepath: str, content: str) -> None:
    gptscript_client = gptscript.GPTScript()
    wksp_file_path = _prepend_base_path("files", filepath)
    await gptscript_client.write_file_in_workspace(
        wksp_file_path, content.encode("utf-8")
    )


async def load_from_gptscript_workspace(filepath: str) -> bytes:
    gptscript_client = gptscript.GPTScript()
    wksp_file_path = _prepend_base_path("files", filepath)
    file_content = await gptscript_client.read_file_in_workspace(wksp_file_path)
    return file_content



def extract_text_from_pdf(pdf_bytes: bytes) -> str:
    """Extracts text from a PDF file given as bytes."""
    pdf_stream = io.BytesIO(pdf_bytes)
    doc = fitz.open(stream=pdf_stream, filetype="pdf")
    
    text = []
    for page in doc:
        text.append(page.get_text("text"))
    
    return "\n".join(text)

def extract_text_from_docx(docx_bytes: bytes) -> str:
    """Extracts text from a Word (.docx) file given as bytes."""
    doc_stream = io.BytesIO(docx_bytes)
    doc = docx.Document(doc_stream)
    
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    """Extracts text from a PowerPoint (.pptx) file given as bytes."""
    ppt_stream = io.BytesIO(pptx_bytes)
    prs = Presentation(ppt_stream)
    
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    return "\n".join(text)